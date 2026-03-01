"""
Prana slide deck generator — matches the app's green/cream palette.
Produces: prana-deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
FOREST      = RGBColor(0x2D, 0x6A, 0x4F)   # #2D6A4F  primary green
MED_GREEN   = RGBColor(0x40, 0x91, 0x6C)   # #40916C
LIGHT_GREEN = RGBColor(0x52, 0xB7, 0x88)   # #52B788
MINT        = RGBColor(0x74, 0xC6, 0x9D)   # #74C69D
PALE_MINT   = RGBColor(0xB7, 0xE4, 0xC7)   # #B7E4C7
CREAM       = RGBColor(0xFA, 0xF6, 0xF0)   # #FAF6F0  app bg
PARCHMENT   = RGBColor(0xE8, 0xDF, 0xD0)   # #E8DFD0  borders
DARK_FOREST = RGBColor(0x1B, 0x43, 0x32)   # #1B4332  primary text
MUTED       = RGBColor(0x6B, 0x7C, 0x6E)   # #6B7C6E  secondary text
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF5, 0xFA, 0xF7)

W = Inches(13.333)   # 16:9 widescreen
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ──────────────────────────────────────────────────────────

def solid_fill(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

def no_line(shape):
    shape.line.fill.background()

def add_rect(slide, l, t, w, h, rgb, line=False):
    s = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    solid_fill(s, rgb)
    if not line:
        no_line(s)
    return s

def add_textbox(slide, l, t, w, h, text, size, bold=False, color=DARK_FOREST,
                align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "SF Pro Display"
    return txb

def accent_bar(slide, top, color=FOREST, h=Inches(0.06)):
    """Thin horizontal rule."""
    add_rect(slide, Inches(0.8), top, Inches(11.73), h, color)

def slide_number(slide, n):
    add_textbox(slide, Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.3),
                str(n), 9, color=MUTED, align=PP_ALIGN.RIGHT)

def logo_bug(slide):
    """Tiny PRANA wordmark top-right."""
    add_textbox(slide, Inches(11.5), Inches(0.22), Inches(1.6), Inches(0.4),
                "PRANA", 11, bold=True, color=FOREST, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Full-bleed cover
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

# Dark forest bg
add_rect(s, 0, 0, W, H, DARK_FOREST)

# Gradient accent band (simulate with two overlapping rects)
add_rect(s, 0, 0, W, Inches(2.5), FOREST)
add_rect(s, 0, 0, W, Inches(1.2), MED_GREEN)

# Decorative circle top-right
c = s.shapes.add_shape(9, Inches(9.8), Inches(-1.5), Inches(5), Inches(5))  # oval
c.fill.solid(); c.fill.fore_color.rgb = MED_GREEN
c.line.fill.background()
c.fill.fore_color.theme_color  # just trigger
from pptx.dml.color import RGBColor as RC
c.fill.fore_color.rgb = RC(0x40, 0x91, 0x6C)

# Small circle bottom-left
c2 = s.shapes.add_shape(9, Inches(-1.5), Inches(5), Inches(4), Inches(4))
c2.fill.solid(); c2.fill.fore_color.rgb = RC(0x2D, 0x6A, 0x4F)
c2.line.fill.background()

# Tagline top
add_textbox(s, Inches(1.2), Inches(0.45), Inches(6), Inches(0.5),
            "BREATHE · TRACK · HEAL", 11, color=PALE_MINT,
            italic=False, bold=False)

# Main title
add_textbox(s, Inches(1.2), Inches(2.2), Inches(8), Inches(1.6),
            "Presentation Title", 54, bold=True, color=WHITE)

# Subtitle
add_textbox(s, Inches(1.2), Inches(4.05), Inches(8), Inches(0.7),
            "Subtitle or presenter name", 22, color=PALE_MINT)

# Date / event line
add_textbox(s, Inches(1.2), Inches(4.85), Inches(8), Inches(0.4),
            "Event  ·  March 2026", 13, color=MINT)

# Thin accent line above title
add_rect(s, Inches(1.2), Inches(2.1), Inches(1.5), Inches(0.05), LIGHT_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Section divider  (full green)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

add_rect(s, 0, 0, W, H, FOREST)

# Left accent panel
add_rect(s, 0, 0, Inches(0.4), H, LIGHT_GREEN)

# Section number
add_textbox(s, Inches(1.2), Inches(2.5), Inches(2), Inches(1),
            "01", 72, bold=True, color=PALE_MINT)

add_rect(s, Inches(1.2), Inches(3.6), Inches(1.2), Inches(0.06), PALE_MINT)

add_textbox(s, Inches(1.2), Inches(3.8), Inches(9), Inches(0.9),
            "Section Title", 40, bold=True, color=WHITE)

add_textbox(s, Inches(1.2), Inches(4.75), Inches(9), Inches(0.5),
            "Brief description of this section", 18, color=PALE_MINT)

slide_number(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Title + content (cream bg)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

add_rect(s, 0, 0, W, H, CREAM)

# Top bar
add_rect(s, 0, 0, W, Inches(1.25), FOREST)

add_textbox(s, Inches(0.8), Inches(0.32), Inches(10), Inches(0.65),
            "Slide Title", 30, bold=True, color=WHITE)

logo_bug(s)

# Body content area placeholder lines (visual guide)
for i, label in enumerate(["Key point one", "Key point two", "Key point three", "Key point four"]):
    y = Inches(1.8 + i * 1.1)
    # bullet dot
    dot = s.shapes.add_shape(9, Inches(0.75), y + Inches(0.12), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = LIGHT_GREEN
    dot.line.fill.background()
    add_textbox(s, Inches(1.15), y, Inches(11.2), Inches(0.5),
                label, 20, color=DARK_FOREST)
    # thin rule
    add_rect(s, Inches(0.8), y + Inches(0.55), Inches(11.73), Inches(0.01), PARCHMENT)

slide_number(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Two-column layout
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

add_rect(s, 0, 0, W, H, CREAM)
add_rect(s, 0, 0, W, Inches(1.25), FOREST)
add_textbox(s, Inches(0.8), Inches(0.32), Inches(10), Inches(0.65),
            "Two-Column Layout", 30, bold=True, color=WHITE)
logo_bug(s)

# Left column card
lc = add_rect(s, Inches(0.7), Inches(1.45), Inches(5.6), Inches(5.5), WHITE)
lc.shadow.inherit = False

# Left column accent top
add_rect(s, Inches(0.7), Inches(1.45), Inches(5.6), Inches(0.12), FOREST)
add_textbox(s, Inches(0.95), Inches(1.65), Inches(5), Inches(0.5),
            "Column One Heading", 16, bold=True, color=DARK_FOREST)
add_rect(s, Inches(0.95), Inches(2.2), Inches(5.1), Inches(0.03), PARCHMENT)
add_textbox(s, Inches(0.95), Inches(2.3), Inches(5.0), Inches(4.0),
            "Body copy or bullet points go here.\n\nAdd your content.", 14, color=MUTED)

# Right column card
rc = add_rect(s, Inches(7.0), Inches(1.45), Inches(5.6), Inches(5.5), WHITE)
rc.shadow.inherit = False
add_rect(s, Inches(7.0), Inches(1.45), Inches(5.6), Inches(0.12), LIGHT_GREEN)
add_textbox(s, Inches(7.25), Inches(1.65), Inches(5), Inches(0.5),
            "Column Two Heading", 16, bold=True, color=DARK_FOREST)
add_rect(s, Inches(7.25), Inches(2.2), Inches(5.1), Inches(0.03), PARCHMENT)
add_textbox(s, Inches(7.25), Inches(2.3), Inches(5.0), Inches(4.0),
            "Body copy or bullet points go here.\n\nAdd your content.", 14, color=MUTED)

slide_number(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Three-stat / icon row
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

add_rect(s, 0, 0, W, H, CREAM)
add_rect(s, 0, 0, W, Inches(1.25), DARK_FOREST)
add_textbox(s, Inches(0.8), Inches(0.32), Inches(10), Inches(0.65),
            "Key Metrics / Statistics", 30, bold=True, color=WHITE)
logo_bug(s)

# Three stat cards
for i, (num, label) in enumerate([("00", "Metric One"), ("00", "Metric Two"), ("00", "Metric Three")]):
    x = Inches(0.7 + i * 4.22)
    card = add_rect(s, x, Inches(1.7), Inches(3.9), Inches(4.8), WHITE)
    # top accent
    add_rect(s, x, Inches(1.7), Inches(3.9), Inches(0.25), [FOREST, LIGHT_GREEN, MINT][i])
    # big number
    add_textbox(s, x + Inches(0.25), Inches(2.2), Inches(3.4), Inches(1.4),
                num, 60, bold=True, color=DARK_FOREST, align=PP_ALIGN.CENTER)
    # label
    add_textbox(s, x + Inches(0.25), Inches(3.7), Inches(3.4), Inches(0.5),
                label, 15, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    # description
    add_textbox(s, x + Inches(0.25), Inches(4.3), Inches(3.4), Inches(1.8),
                "Short description of\nthis metric", 12, color=MUTED, align=PP_ALIGN.CENTER)

slide_number(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Quote / highlight
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

add_rect(s, 0, 0, W, H, OFF_WHITE)

# Big green left bar
add_rect(s, 0, 0, Inches(0.55), H, FOREST)

# Quote mark decorative
add_textbox(s, Inches(1.3), Inches(0.8), Inches(2), Inches(2),
            "\u201C", 120, bold=True, color=PALE_MINT)

add_textbox(s, Inches(1.5), Inches(2.0), Inches(10.0), Inches(2.4),
            "Insert a meaningful quote, insight, or\nhighlighted statement here.",
            32, bold=False, color=DARK_FOREST)

accent_bar(s, Inches(4.55), LIGHT_GREEN, Inches(0.05))

add_textbox(s, Inches(1.5), Inches(4.7), Inches(6), Inches(0.4),
            "— Attribution or source", 14, italic=True, color=MUTED)

logo_bug(s)
slide_number(s, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Image placeholder + caption
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

add_rect(s, 0, 0, W, H, CREAM)
add_rect(s, 0, 0, W, Inches(1.25), FOREST)
add_textbox(s, Inches(0.8), Inches(0.32), Inches(10), Inches(0.65),
            "Visual / Demo", 30, bold=True, color=WHITE)
logo_bug(s)

# Left: image placeholder box
img_box = add_rect(s, Inches(0.7), Inches(1.5), Inches(7.8), Inches(5.5), PARCHMENT)
# dashed border effect — thin inner rect
add_rect(s, Inches(0.75), Inches(1.55), Inches(7.7), Inches(5.4), PARCHMENT)
add_textbox(s, Inches(0.7), Inches(3.5), Inches(7.8), Inches(1.0),
            "[ Image / Screenshot / Chart ]", 18, color=MUTED, align=PP_ALIGN.CENTER)

# Right: caption / description
add_textbox(s, Inches(9.0), Inches(1.8), Inches(3.8), Inches(0.6),
            "Caption", 20, bold=True, color=DARK_FOREST)
accent_bar(s, Inches(2.5), LIGHT_GREEN, Inches(0.04))
add_textbox(s, Inches(9.0), Inches(2.65), Inches(3.8), Inches(3.5),
            "Describe the image, chart, or demo shown on the left. Use 2–4 sentences.", 14, color=MUTED)

slide_number(s, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Team / people grid (3-up)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

add_rect(s, 0, 0, W, H, CREAM)
add_rect(s, 0, 0, W, Inches(1.25), DARK_FOREST)
add_textbox(s, Inches(0.8), Inches(0.32), Inches(10), Inches(0.65),
            "Meet the Team", 30, bold=True, color=WHITE)
logo_bug(s)

for i in range(3):
    x = Inches(0.85 + i * 4.17)
    # Avatar circle
    av = s.shapes.add_shape(9, x + Inches(0.8), Inches(1.6), Inches(2.4), Inches(2.4))
    av.fill.solid()
    av.fill.fore_color.rgb = [PALE_MINT, PARCHMENT, PALE_MINT][i]
    av.line.fill.background()
    # Name
    add_textbox(s, x, Inches(4.2), Inches(4.0), Inches(0.5),
                "Name Surname", 16, bold=True, color=DARK_FOREST, align=PP_ALIGN.CENTER)
    # Role
    add_textbox(s, x, Inches(4.75), Inches(4.0), Inches(0.4),
                "Role / Title", 13, color=MUTED, align=PP_ALIGN.CENTER)
    # Bio line
    add_textbox(s, x, Inches(5.2), Inches(4.0), Inches(1.5),
                "Short bio or area of focus", 11, color=MUTED, align=PP_ALIGN.CENTER)

slide_number(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Timeline / roadmap
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

add_rect(s, 0, 0, W, H, CREAM)
add_rect(s, 0, 0, W, Inches(1.25), FOREST)
add_textbox(s, Inches(0.8), Inches(0.32), Inches(10), Inches(0.65),
            "Roadmap / Timeline", 30, bold=True, color=WHITE)
logo_bug(s)

# Horizontal spine
add_rect(s, Inches(0.8), Inches(3.9), Inches(11.73), Inches(0.08), LIGHT_GREEN)

phases = ["Phase 1", "Phase 2", "Phase 3", "Phase 4"]
colors = [FOREST, LIGHT_GREEN, MINT, PALE_MINT]
tcols  = [WHITE, WHITE, DARK_FOREST, DARK_FOREST]
for i, (ph, col, tc) in enumerate(zip(phases, colors, tcols)):
    x = Inches(0.8 + i * 3.0)
    # Node circle
    node = s.shapes.add_shape(9, x + Inches(0.85), Inches(3.55), Inches(0.78), Inches(0.78))
    node.fill.solid(); node.fill.fore_color.rgb = col
    node.line.fill.background()
    # Phase label (above)
    add_textbox(s, x, Inches(2.6), Inches(2.5), Inches(0.5),
                ph, 15, bold=True, color=DARK_FOREST, align=PP_ALIGN.CENTER)
    add_textbox(s, x, Inches(3.0), Inches(2.5), Inches(0.4),
                "Q1 2026", 11, color=MUTED, align=PP_ALIGN.CENTER)
    # Description (below)
    add_textbox(s, x, Inches(4.55), Inches(2.5), Inches(2.0),
                "Description of deliverables or milestones.", 11, color=MUTED, align=PP_ALIGN.CENTER)

slide_number(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Closing / Thank you
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

add_rect(s, 0, 0, W, H, DARK_FOREST)

# decorative circles
for (cx, cy, sz, col) in [
    (Inches(10.5), Inches(-1), Inches(5), MED_GREEN),
    (Inches(-1),   Inches(5),  Inches(4), FOREST),
    (Inches(11),   Inches(5.5),Inches(3), LIGHT_GREEN),
]:
    c = s.shapes.add_shape(9, cx, cy, sz, sz)
    c.fill.solid(); c.fill.fore_color.rgb = col
    c.line.fill.background()

add_rect(s, Inches(1.2), Inches(3.4), Inches(1.5), Inches(0.05), LIGHT_GREEN)

add_textbox(s, Inches(1.2), Inches(1.8), Inches(9), Inches(1.8),
            "Thank You", 64, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_textbox(s, Inches(1.2), Inches(3.6), Inches(8), Inches(0.6),
            "contact@example.com  ·  prana.app", 16, color=PALE_MINT)

add_textbox(s, Inches(1.2), Inches(4.4), Inches(4), Inches(0.4),
            "PRANA", 13, bold=True, color=MINT)

add_textbox(s, Inches(1.2), Inches(4.9), Inches(6), Inches(0.4),
            "BREATHE · TRACK · HEAL", 10, color=MUTED)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/harshithreddy/Desktop/Hackathon/prana-deck.pptx"
prs.save(out)
print(f"Saved → {out}")
